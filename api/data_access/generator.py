import os
import requests
from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pathlib import Path
from openai import OpenAI

load_dotenv()
client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))


def generar_contenido(prompt):
    prompt = prompt.strip()
    response = client.chat.completions.create(
        model="gpt-4",
        messages=[
            {"role": "system", "content": (
                "Eres un generador profesional de presentaciones de PowerPoint. "
                "Cada diapositiva debe tener:\n"
                "- Un título claro y breve.\n"
                "- Un párrafo explicativo de 2 a 4 líneas.\n"
                "- Una imagen relevante dentro de los tags [image: descripción visual]."
            )},
            {"role": "user", "content": prompt}
        ]
    )
    texto = response.choices[0].message.content
    slides = []

    for bloque in texto.split("\n\n"):
        slide = {"titulo": "", "contenido": "", "imagen": ""}
        for linea in bloque.split("\n"):
            if linea.lower().startswith("título:"):
                slide["titulo"] = linea.split(":", 1)[1].strip()
            elif "[image:" in linea:
                slide["imagen"] = linea.split("[image:", 1)[1].split("]")[0].strip()
            else:
                slide["contenido"] += linea.replace("Párrafo:", "").strip() + " "
        if slide["titulo"] or slide["contenido"]:
            slides.append(slide)

    return slides


def generar_imagen(prompt_img, filename):
    try:
        image_response = client.images.generate(
            model="dall-e-3",
            prompt=prompt_img,
            size="1024x1024",
            quality="standard",
            n=1,
        )
        url = image_response.data[0].url
        img_data = requests.get(url).content
        with open(filename, 'wb') as f:
            f.write(img_data)
        return True
    except Exception as e:
        print("❌ Error al generar imagen:", e)
        return False


def aplicar_fondo_claro(slide, color=RGBColor(250, 250, 250)):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def agregar_titulo(slide, texto, color=RGBColor(52, 73, 94)):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8.5), Inches(1))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = texto
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.LEFT


def agregar_caja_contenido(slide, texto, color=RGBColor(80, 80, 80)):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(8.5), Inches(2.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = texto
    p.font.size = Pt(20)
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.LEFT


def agregar_caja_decorativa(slide):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(1.2), Inches(9), Inches(3)
    )
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    line = shape.line
    line.color.rgb = RGBColor(200, 200, 200)
    shape.shadow.inherit = False
    shape.shadow.blur_radius = 5.0


def insertar_imagen(slide, img_path):
    slide.shapes.add_picture(str(img_path), Inches(5.5), Inches(2.8), height=Inches(3.5))


def generar_pptx(slides, nombre):
    folder = Path("media/presentaciones")
    folder.mkdir(parents=True, exist_ok=True)
    ruta = folder / f"{nombre}.pptx"

    prs = Presentation()
    blank_layout = prs.slide_layouts[6]

    # Slide de título inicial
    slide_title = prs.slides.add_slide(prs.slide_layouts[0])
    slide_title.shapes.title.text = nombre.replace("_", " ").capitalize()
    subtitle = slide_title.placeholders[1]
    subtitle.text = "Presentación generada automáticamente"
    subtitle.text_frame.paragraphs[0].font.size = Pt(18)

    # Diapositivas del contenido
    for i, slide_data in enumerate(slides):
        slide = prs.slides.add_slide(blank_layout)
        aplicar_fondo_claro(slide)

        agregar_caja_decorativa(slide)
        agregar_titulo(slide, slide_data["titulo"])
        agregar_caja_contenido(slide, slide_data["contenido"])

        if slide_data["imagen"]:
            img_path = folder / f"{nombre}_img{i+1}.png"
            if generar_imagen(slide_data["imagen"], img_path):
                insertar_imagen(slide, img_path)

        slide.notes_slide.notes_text_frame.text = "Transición recomendada: fundido"

    prs.save(ruta)
    return str(ruta)
